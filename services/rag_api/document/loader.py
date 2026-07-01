from __future__ import annotations

import csv
import io
from pathlib import Path
from collections.abc import Iterable

from charset_normalizer import from_bytes
from docx import Document

from services.rag_api.document.cleaner import clean_text
from services.rag_api.exceptions import DOC_LOAD_ERROR_MESSAGE, DocumentLoadError


SUPPORTED_SUFFIXES = {".docx", ".txt", ".pdf", ".xlsx", ".xlsm", ".csv", ".pptx"}


def read_txt_file(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gbk", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    match = from_bytes(data).best()
    if match is None:
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE)
    return str(match)


def read_docx_file(path: Path) -> str:
    document = Document(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    table_cells: list[str] = []
    for table in document.tables:
        for row in table.rows:
            table_cells.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
    return "\n".join(paragraphs + table_cells)


def read_pdf_file(path: Path) -> str:
    pymupdf = _load_pymupdf()
    document = None
    try:
        document = pymupdf.open(str(path))
        text = "\n".join((page.get_text("text") or "").strip() for page in document)
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc
    finally:
        if document is not None:
            document.close()
    if not text.strip():
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE)
    return text


def read_csv_file(path: Path) -> str:
    text = read_txt_file(path)
    rows: list[str] = []
    for row in csv.reader(io.StringIO(text)):
        values = [cell.strip() for cell in row if cell and cell.strip()]
        if values:
            rows.append(" | ".join(values))
    return "\n".join(rows)


def read_excel_file(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc

    workbook = None
    try:
        workbook = load_workbook(path, read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in workbook.worksheets:
            sheet_rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    sheet_rows.append(" | ".join(values))
            if sheet_rows:
                rows.append(f"工作表：{sheet.title}")
                rows.extend(sheet_rows)
        return "\n".join(rows)
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc
    finally:
        if workbook is not None:
            workbook.close()


def read_pptx_file(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc

    try:
        presentation = Presentation(str(path))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_parts.append(text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if values:
                            slide_parts.append(" | ".join(values))
            if slide_parts:
                parts.append(f"幻灯片 {slide_index}")
                parts.extend(slide_parts)
        return "\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc


def _load_pymupdf():
    try:
        import pymupdf  # type: ignore

        return pymupdf
    except Exception:
        pass
    try:
        import fitz  # type: ignore

        if not hasattr(fitz, "open"):
            raise RuntimeError("imported fitz is not PyMuPDF")
        return fitz
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc


def load_documents(kb_dir: Path | str | Iterable[Path | str]) -> list[dict]:
    dirs = _normalize_dirs(kb_dir)
    files = _scan_supported_files(dirs)
    if not files:
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE)
    documents = [doc for doc in (load_document(path) for path in files) if doc.get("content")]
    if not documents:
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE)
    return documents


def load_document(path: Path | str) -> dict:
    file_path = Path(path)
    try:
        suffix = file_path.suffix.lower()
        if suffix == ".docx":
            raw = read_docx_file(file_path)
        elif suffix == ".pdf":
            raw = read_pdf_file(file_path)
        elif suffix in {".xlsx", ".xlsm"}:
            raw = read_excel_file(file_path)
        elif suffix == ".csv":
            raw = read_csv_file(file_path)
        elif suffix == ".pptx":
            raw = read_pptx_file(file_path)
        else:
            raw = read_txt_file(file_path)
        content = clean_text(raw)
        return {"source_file": file_path.name, "source_path": str(file_path), "content": content}
    except Exception as exc:  # noqa: BLE001
        raise DocumentLoadError(DOC_LOAD_ERROR_MESSAGE) from exc


def has_supported_documents(kb_dir: Path | str | Iterable[Path | str]) -> bool:
    return bool(_scan_supported_files(_normalize_dirs(kb_dir)))


def scan_supported_files(kb_dir: Path | str | Iterable[Path | str]) -> list[Path]:
    return _scan_supported_files(_normalize_dirs(kb_dir))


def _normalize_dirs(kb_dir: Path | str | Iterable[Path | str]) -> list[Path]:
    if isinstance(kb_dir, (str, Path)):
        raw_dirs = [kb_dir]
    else:
        raw_dirs = list(kb_dir)
    result: list[Path] = []
    for item in raw_dirs:
        path = Path(item).expanduser().resolve()
        if path not in result:
            result.append(path)
    return result


def _scan_supported_files(dirs: list[Path]) -> list[Path]:
    files: list[Path] = []
    for directory in dirs:
        if not directory.exists() or not directory.is_dir():
            continue
        files.extend(path for path in directory.rglob("*") if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES)
    return sorted(files, key=lambda path: str(path).lower())
